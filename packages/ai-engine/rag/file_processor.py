import io
import re
from typing import Dict, Any

class FileProcessor:
    """
    Highly resilient text extractor for PDF, DOCX, PPTX, and standard text files.
    Features robust error handling and low-level regex fallback to prevent pipeline crashes.
    """
    @staticmethod
    def extract_text(file_bytes: bytes, file_type: str) -> str:
        if not file_bytes:
            return ""

        file_type = file_type.lower().strip().replace(".", "")
        text = ""

        try:
            if file_type == "pdf":
                text = FileProcessor._extract_pdf(file_bytes)
            elif file_type in ["doc", "docx"]:
                text = FileProcessor._extract_docx(file_bytes)
            elif file_type in ["ppt", "pptx"]:
                text = FileProcessor._extract_pptx(file_bytes)
            elif file_type in ["txt", "md", "csv", "json"]:
                text = file_bytes.decode("utf-8", errors="ignore")
            else:
                # Direct decoding attempt for unknown formats
                text = file_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            print(f"[WARN] Standard parser failed for {file_type}: {e}. Running structural regex fallback.")
            text = FileProcessor._extract_regex_fallback(file_bytes)

        # If standard parser succeeded but returned zero text, try fallback
        if not text.strip():
            text = FileProcessor._extract_regex_fallback(file_bytes)

        return text

    @staticmethod
    def _extract_pdf(file_bytes: bytes) -> str:
        import pypdf
        text_list = []
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_list.append(page_text)
        return "\n".join(text_list)

    @staticmethod
    def _extract_docx(file_bytes: bytes) -> str:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        text_list = []
        
        # Extract from paragraphs
        for p in doc.paragraphs:
            if p.text.strip():
                text_list.append(p.text)
                
        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_list.append(" | ".join(row_text))
                    
        return "\n".join(text_list)

    @staticmethod
    def _extract_pptx(file_bytes: bytes) -> str:
        import pptx
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        text_list = []
        
        for i, slide in enumerate(prs.slides):
            text_list.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_list.append(shape.text)
        return "\n".join(text_list)

    @staticmethod
    def _extract_regex_fallback(file_bytes: bytes) -> str:
        """
        Scans binary data for continuous alphanumeric strings representing text blocks.
        Acts as a fail-safe fallback mechanism.
        """
        # Find continuous strings of readable characters (letters, numbers, common punctuation, spaces)
        pattern = re.compile(rb'[a-zA-Z0-9\s\.,;:!\?\-\(\)\[\]"\'\/\\@#\$%\^\&\*\+\=\_]{6,}')
        matches = pattern.findall(file_bytes)
        
        extracted_lines = []
        for match in matches:
            try:
                line = match.decode("utf-8", errors="ignore").strip()
                # Exclude strings with too many non-letters or weird spacing
                if len(line) > 8 and not line.startswith("http") and not line.replace(" ", "").isalnum() is False:
                    extracted_lines.append(line)
            except Exception:
                continue
                
        return "\n".join(extracted_lines)
